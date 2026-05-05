import os
import sys
import json
import shutil
import hashlib
import mimetypes
import traceback
import subprocess
import threading
from pathlib import Path
from typing import Optional

# ── Text extraction helpers ──────────────────────────────────────────────────

def extract_text_from_file(filepath: str, max_chars: int = 2000) -> str:
    """Try to extract readable text from any file. Returns '' on failure."""
    path = Path(filepath)
    suffix = path.suffix.lower()

    # Plain text / code
    text_exts = {
        ".txt", ".md", ".csv", ".json", ".xml", ".html", ".htm",
        ".py", ".js", ".ts", ".java", ".c", ".cpp", ".cs", ".go",
        ".rs", ".rb", ".php", ".sh", ".yaml", ".yml", ".toml", ".ini",
        ".log", ".rst", ".tex", ".sql", ".r", ".swift", ".kt"
    }
    if suffix in text_exts:
        try:
            with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
                return f.read(max_chars)
        except Exception:
            return ""

    # PDF
    if suffix == ".pdf":
        return _extract_pdf(filepath, max_chars)

    # Word documents
    if suffix in (".docx", ".doc"):
        return _extract_docx(filepath, max_chars)

    # Excel / spreadsheets
    if suffix in (".xlsx", ".xls", ".ods"):
        return _extract_excel(filepath, max_chars)

    # PowerPoint
    if suffix in (".pptx", ".ppt"):
        return _extract_pptx(filepath, max_chars)

    # Images — return metadata only
    if suffix in (".jpg", ".jpeg", ".png", ".gif", ".bmp", ".webp", ".svg", ".tiff"):
        return f"[Image file: {path.name}]"

    # Audio/video — metadata only
    if suffix in (".mp3", ".mp4", ".wav", ".mkv", ".avi", ".mov", ".flac", ".ogg"):
        return f"[Media file: {path.name}]"

    # Zip — list contents
    if suffix == ".zip":
        return _list_zip(filepath, max_chars)

    # Fallback: try reading as text
    try:
        with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
            return f.read(max_chars)
    except Exception:
        return ""


def _extract_pdf(filepath: str, max_chars: int) -> str:
    try:
        import fitz  # PyMuPDF
        doc = fitz.open(filepath)
        text = ""
        for page_num in range(min(3, len(doc))):  # First 3 pages max
            text += doc[page_num].get_text()
            if len(text) >= max_chars:
                break
        doc.close()
        return text[:max_chars]
    except ImportError:
        pass
    # Fallback: pdfminer
    try:
        from pdfminer.high_level import extract_text
        text = extract_text(filepath, maxpages=3)
        return text[:max_chars]
    except Exception:
        return "[PDF file - install PyMuPDF for text extraction]"


def _extract_docx(filepath: str, max_chars: int) -> str:
    try:
        import docx
        doc = docx.Document(filepath)
        text = "\n".join(p.text for p in doc.paragraphs)
        return text[:max_chars]
    except Exception:
        return "[Word document - install python-docx for text extraction]"


def _extract_excel(filepath: str, max_chars: int) -> str:
    try:
        import openpyxl
        wb = openpyxl.load_workbook(filepath, read_only=True, data_only=True)
        text = ""
        for sheet in wb.sheetnames[:2]:
            ws = wb[sheet]
            for row in ws.iter_rows(max_row=20, values_only=True):
                row_text = " | ".join(str(c) for c in row if c is not None)
                if row_text.strip():
                    text += row_text + "\n"
            if len(text) >= max_chars:
                break
        return text[:max_chars]
    except Exception:
        return "[Spreadsheet file]"


def _extract_pptx(filepath: str, max_chars: int) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(filepath)
        text = ""
        for slide in prs.slides[:5]:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
            if len(text) >= max_chars:
                break
        return text[:max_chars]
    except Exception:
        return "[PowerPoint file]"


def _list_zip(filepath: str, max_chars: int) -> str:
    try:
        import zipfile
        with zipfile.ZipFile(filepath) as zf:
            names = zf.namelist()
        return "ZIP archive containing: " + ", ".join(names[:30])
    except Exception:
        return "[ZIP archive]"


# ── Ollama / Gemma 2B classifier ─────────────────────────────────────────────

CLASSIFICATION_PROMPT = """You are a file classification expert. Analyze the following batch of files and classify EACH of them into a folder hierarchy.

Rules:
- Create 1-3 levels of folder hierarchy (e.g., "History/Ancient/Wars")
- Be specific but not too granular
- Use common sense category names
- Return ONLY a JSON array of objects, one for each file in the EXACT same order.
- Each object MUST have these fields:
  - "filename": the exact file name
  - "category": the top-level folder name
  - "subcategory": a more specific subfolder (can be null)
  - "detail": even more specific subfolder (can be null, only if very clear)
  - "reason": one sentence explaining why

Examples of good hierarchies:
- History/Ancient/Wars
- Programming/Python/Web Development
- Medical/Cardiology/Research Papers  
- Finance/Tax Documents/2024

Files to classify:
{files_content}

Respond ONLY with a valid JSON array, no markdown, no explanation:"""


def classify_batch_with_ollama(
    files_info: list[dict],
    model: str = "gemma2:2b",
    ollama_url: str = "http://localhost:11434"
) -> list[dict]:
    """Send content to Ollama and get classifications back for a batch."""
    import urllib.request

    files_text = ""
    for idx, info in enumerate(files_info):
        content = info["content"][:1000] if info["content"] else "(no text content)"
        files_text += f"File {idx+1}:\nName: {info['filename']}\nType: {info['filetype']}\nContent preview:\n---\n{content}\n---\n\n"

    prompt = CLASSIFICATION_PROMPT.format(files_content=files_text.strip())

    payload = json.dumps({
        "model": model,
        "prompt": prompt,
        "stream": False,
        "options": {
            "temperature": 0.1,   # Low temp for consistent classification
            "num_predict": 1024   # Larger to accommodate JSON array
        }
    }).encode("utf-8")

    req = urllib.request.Request(
        f"{ollama_url}/api/generate",
        data=payload,
        headers={"Content-Type": "application/json"},
        method="POST"
    )

    try:
        with urllib.request.urlopen(req, timeout=120) as resp:
            result = json.loads(resp.read().decode("utf-8"))
            raw_text = result.get("response", "")
    except Exception as e:
        raise RuntimeError(f"Ollama request failed: {e}")

    # Parse JSON from response
    try:
        clean = raw_text.strip()
        if clean.startswith("```"):
            clean = clean.split("```")[1]
            if clean.startswith("json"):
                clean = clean[4:]
        parsed = json.loads(clean.strip())
        if isinstance(parsed, list):
            return parsed
        # If model returned a single object instead of a list
        return [parsed]
    except json.JSONDecodeError:
        # Fallback regex extraction
        import re
        match = re.search(r'\[.*\]', raw_text, re.DOTALL)
        if match:
            try:
                parsed = json.loads(match.group())
                if isinstance(parsed, list):
                    return parsed
            except Exception:
                pass
        return []


# ── Folder path builder ───────────────────────────────────────────────────────

def build_dest_path(output_dir: str, classification: dict, filename: str) -> Path:
    """Build the destination path from classification result."""
    parts = [output_dir]

    category = classification.get("category", "Uncategorized").strip()
    subcategory = classification.get("subcategory")
    detail = classification.get("detail")

    parts.append(_sanitize(category))
    if subcategory and subcategory.strip():
        parts.append(_sanitize(subcategory.strip()))
    if detail and detail.strip():
        parts.append(_sanitize(detail.strip()))

    dest_folder = Path(*parts)
    dest_folder.mkdir(parents=True, exist_ok=True)

    # Handle filename collision
    dest_file = dest_folder / filename
    if dest_file.exists():
        stem = Path(filename).stem
        ext = Path(filename).suffix
        h = hashlib.md5(filename.encode()).hexdigest()[:6]
        dest_file = dest_folder / f"{stem}_{h}{ext}"

    return dest_file


def _sanitize(name: str) -> str:
    """Remove characters not safe for folder names."""
    import re
    name = re.sub(r'[<>:"/\\|?*]', '_', name)
    return name.strip(" .")[:80] or "Misc"


# ── Main sorting function ─────────────────────────────────────────────────────

def _sort_batch(args):
    """Process a batch of files."""
    batch_indices, filepaths, total, output_dir, model, ollama_url, copy_mode, progress_callback, results_lock, results = args
    
    files_info = []
    batch_results = []
    
    for i, filepath in zip(batch_indices, filepaths):
        filename = os.path.basename(filepath)
        result = {
            "file": filepath,
            "filename": filename,
            "dest": None,
            "category": None,
            "subcategory": None,
            "detail": None,
            "reason": None,
            "error": None,
            "index": i
        }
        batch_results.append(result)

        if progress_callback:
            progress_callback(i, total, filename, "reading", None)

        try:
            filetype = mimetypes.guess_type(filepath)[0] or Path(filepath).suffix
            content = extract_text_from_file(filepath, max_chars=1000)
            files_info.append({
                "filename": filename,
                "filetype": filetype,
                "content": content,
                "result_ref": result
            })
        except Exception as e:
            result["error"] = f"Reading error: {e}"
            traceback.print_exc()

    # Filter out files that had reading errors
    valid_files_info = [info for info in files_info if not info["result_ref"].get("error")]

    if valid_files_info:
        for info in valid_files_info:
            if progress_callback:
                progress_callback(info["result_ref"]["index"], total, info["filename"], "classifying", None)

        try:
            classifications = classify_batch_with_ollama(valid_files_info, model, ollama_url)
            
            # Match classifications back to valid_files_info based on list position or filename
            for idx, info in enumerate(valid_files_info):
                classification = {}
                # Try to find matching filename first
                matched = False
                for c in classifications:
                    if isinstance(c, dict) and c.get("filename") == info["filename"]:
                        classification = c
                        matched = True
                        break
                
                # Fallback to positional mapping if array length matches
                if not matched and idx < len(classifications):
                    c = classifications[idx]
                    if isinstance(c, dict):
                        classification = c
                
                if not classification:
                    classification = {
                        "category": "Uncategorized",
                        "reason": "Failed to parse from batch response"
                    }

                result = info["result_ref"]
                filepath = result["file"]
                filename = result["filename"]

                with results_lock:
                    dest = build_dest_path(output_dir, classification, filename)

                try:
                    if copy_mode:
                        shutil.copy2(filepath, dest)
                    else:
                        shutil.move(filepath, dest)

                    result.update({
                        "dest": str(dest),
                        "category": classification.get("category"),
                        "subcategory": classification.get("subcategory"),
                        "detail": classification.get("detail"),
                        "reason": classification.get("reason"),
                    })

                    if progress_callback:
                        progress_callback(result["index"], total, filename, "done", result)
                except Exception as e:
                    result["error"] = f"File operation error: {e}"
                    if progress_callback:
                        progress_callback(result["index"], total, filename, "error", result)

        except Exception as e:
            error_msg = f"Classification batch error: {e}"
            traceback.print_exc()
            for info in valid_files_info:
                info["result_ref"]["error"] = error_msg
                if progress_callback:
                    progress_callback(info["result_ref"]["index"], total, info["filename"], "error", info["result_ref"])

    with results_lock:
        results.extend(batch_results)

    return batch_results


def sort_files(
    filepaths: list[str],
    output_dir: str,
    model: str = "gemma2:2b",
    ollama_url: str = "http://localhost:11434",
    copy_mode: bool = True,
    progress_callback=None,
    workers: int = 2
) -> list[dict]:
    """
    Sort files into categorized folders in parallel batches.
    workers=2 with batch_size=5 means up to 10 files processed concurrently.
    """
    from concurrent.futures import ThreadPoolExecutor, as_completed

    total = len(filepaths)
    results = []
    results_lock = threading.Lock()
    
    batch_size = 5
    batches = []
    
    for i in range(0, total, batch_size):
        batch_paths = filepaths[i:i+batch_size]
        batch_indices = list(range(i, i + len(batch_paths)))
        batches.append((
            batch_indices, batch_paths, total, output_dir, model, ollama_url, copy_mode,
            progress_callback, results_lock, results
        ))

    with ThreadPoolExecutor(max_workers=workers) as executor:
        futures = [executor.submit(_sort_batch, args) for args in batches]
        for future in as_completed(futures):
            try:
                future.result()
            except Exception as e:
                traceback.print_exc()

    results.sort(key=lambda r: next(
        (i for i, fp in enumerate(filepaths) if fp == r["file"]), 0
    ))
    return results


# ── Folder scanner ────────────────────────────────────────────────────────────

def scan_folder(folder_path: str) -> list[str]:
    """Return all files in a folder (non-recursive)."""
    folder = Path(folder_path)
    if not folder.exists():
        raise ValueError(f"Folder not found: {folder_path}")
    if not folder.is_dir():
        raise ValueError(f"Not a folder: {folder_path}")
    return [str(f) for f in folder.iterdir() if f.is_file()]


# ── CLI entrypoint ────────────────────────────────────────────────────────────

if __name__ == "__main__":
    import argparse

    parser = argparse.ArgumentParser(description="COBBLE - AI File Sorter using Ollama")
    parser.add_argument("input", nargs="+", help="Files or a folder path to sort")
    parser.add_argument("-o", "--output", required=True, help="Output directory")
    parser.add_argument("--model", default="gemma2:2b", help="Ollama model")
    parser.add_argument("--url", default="http://localhost:11434", help="Ollama URL")
    parser.add_argument("--move", action="store_true", help="Move files instead of copying")
    args = parser.parse_args()

    # If a single folder path is given, scan it
    filepaths = []
    if len(args.input) == 1 and os.path.isdir(args.input[0]):
        print(f"📂 Scanning folder: {args.input[0]}")
        filepaths = scan_folder(args.input[0])
        print(f"   Found {len(filepaths)} files\n")
    else:
        filepaths = args.input

    def cb(i, total, filename, status, result):
        icons = {"reading": "📄", "classifying": "🤖", "done": "✅", "error": "❌"}
        icon = icons.get(status, "⏳")
        print(f"[{i+1}/{total}] {icon} {filename} — {status}", flush=True)
        if result and result.get("category"):
            path = " / ".join(filter(None, [
                result.get("category"),
                result.get("subcategory"),
                result.get("detail")
            ]))
            print(f"         → {path}", flush=True)

    results = sort_files(
        filepaths=filepaths,
        output_dir=args.output,
        model=args.model,
        ollama_url=args.url,
        copy_mode=not args.move,
        progress_callback=cb
    )

    print(f"\n{'='*50}")
    print(f"Done! {sum(1 for r in results if not r['error'])} / {len(results)} files sorted")
    errors = [r for r in results if r["error"]]
    if errors:
        print(f"\n{len(errors)} error(s):")
        for r in errors:
            print(f"  ✗ {r['filename']}: {r['error']}")
